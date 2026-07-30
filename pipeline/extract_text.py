"""Low-level text extraction: PDF (PyMuPDF), PPTX (python-pptx), XLSX helpers."""
from pathlib import Path
import hashlib


def file_hash(path: Path) -> str:
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(1 << 20), b""):
            h.update(chunk)
    return h.hexdigest()


def pdf_text(path: Path) -> str:
    import fitz  # PyMuPDF

    doc = fitz.open(path)
    try:
        return "\n".join(page.get_text() for page in doc)
    finally:
        doc.close()


def pptx_text(path: Path) -> tuple[str, bool]:
    """Returns (extracted_text, is_image_only).

    Many investor decks flatten each slide into a single full-slide picture;
    in that case there is no text run to extract and is_image_only is True so
    callers can skip KPI parsing instead of silently returning nothing.
    """
    from pptx import Presentation

    prs = Presentation(str(path))
    chunks: list[str] = []
    any_text = False
    any_picture = False
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # PICTURE
                any_picture = True
            if shape.has_text_frame and shape.text_frame.text.strip():
                chunks.append(shape.text_frame.text)
                any_text = True
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            chunks.append(slide.notes_slide.notes_text_frame.text)
            any_text = True
    is_image_only = any_picture and not any_text
    return "\n".join(chunks), is_image_only
